import io
import os
import sys

_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

"""
📈 Visualization Studio — Consolidated Visualization & Dashboard Hub (Premium)
Advanced Visuals, Executive Dashboard Builder, Presentation Deck Generator,
Chart Data Extractor, and an intelligent Auto-Recommendation engine.
"""

import numpy as np
import pandas as pd
import streamlit as st

from modules.page_bootstrap import render_standard_footer, setup_page
from modules.session_manager import get_active_dataframe
from modules.shared_ui import (
    hero_card,
    render_dataset_context_banner,
    render_export_buttons,
    section_header,
)

try:
    import plotly.express as px
    import plotly.graph_objects as go

    PLOTLY_AVAILABLE = True
except ImportError:
    PLOTLY_AVAILABLE = False

try:
    import statsmodels.api as _sm_check  # noqa: F401

    STATSMODELS_AVAILABLE = True
except ImportError:
    STATSMODELS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import kaleido  # noqa: F401

    KALEIDO_AVAILABLE = True
except ImportError:
    KALEIDO_AVAILABLE = False


def get_df() -> pd.DataFrame:
    """Retrieve active session dataframe or supply a clean mock dataset."""
    df = get_active_dataframe()
    if df is None or df.empty:
        np.random.seed(42)
        return pd.DataFrame(
            {
                "Category": np.random.choice(
                    ["Type A", "Type B", "Type C", "Type D"], 120
                ),
                "SubCategory": np.random.choice(
                    ["North-East", "South-West", "Central"], 120
                ),
                "Value_A": np.random.normal(55, 12, 120),
                "Value_B": np.random.normal(32, 9, 120),
                "Metric": np.random.uniform(5, 95, 120),
                "Date": pd.date_range(start="2026-01-01", periods=120, freq="D"),
            }
        )
    return df


def build_chart(
    chart_type: str,
    df: pd.DataFrame,
    x: str = None,
    y: str = None,
    color: str = None,
    facet: str = None,
    size: str = None,
    height: int = 420,
    **kwargs,
) -> go.Figure:
    """Robust factory function to construct Plotly charts safely without type-mismatch failures."""
    if not PLOTLY_AVAILABLE or df is None or df.empty:
        return None

    try:
        template = "plotly_dark"
        working_df = df.copy()

        x = x or kwargs.get("x")
        y = y or kwargs.get("y")
        color = color or kwargs.get("color")

        num_cols = working_df.select_dtypes(include=[np.number]).columns.tolist()
        cat_cols = working_df.select_dtypes(
            include=["object", "category"]
        ).columns.tolist()
        all_cols = working_df.columns.tolist()

        fallback_x = (
            x
            if x and x in all_cols
            else (cat_cols[0] if cat_cols else all_cols[0])
        )
        fallback_y = (
            y if y and y in all_cols else (num_cols[0] if num_cols else None)
        )

        # Enforce numeric types safely on chart types that require numerical Y values
        continuous_charts = [
            "Scatter Plot",
            "Line Chart",
            "Area Chart",
            "Funnel Chart",
        ]
        if chart_type in continuous_charts and fallback_y:
            if fallback_y not in num_cols:
                converted_y = pd.to_numeric(working_df[fallback_y], errors="coerce")
                if not converted_y.isna().all():
                    working_df[fallback_y] = converted_y
                elif num_cols:
                    fallback_y = num_cols[0]

        valid_color = color if color and color in all_cols else None
        valid_size = size if size and size in num_cols else None

        if chart_type == "Histogram":
            fig = px.histogram(
                working_df,
                x=fallback_x,
                color=valid_color,
                barmode="group",
                template=template,
                height=height,
            )
        elif chart_type == "Scatter Plot":
            use_trendline = (
                "ols"
                if (
                    len(working_df) > 5
                    and fallback_y
                    and STATSMODELS_AVAILABLE
                    and fallback_x in num_cols
                )
                else None
            )
            fig = px.scatter(
                working_df,
                x=fallback_x,
                y=fallback_y,
                color=valid_color,
                size=valid_size,
                template=template,
                height=height,
                trendline=use_trendline,
            )
        elif chart_type == "Bar Chart":
            fig = px.bar(
                working_df,
                x=fallback_x,
                y=fallback_y,
                color=valid_color,
                barmode="group",
                template=template,
                height=height,
            )
        elif chart_type == "Line Chart":
            fig = px.line(
                working_df,
                x=fallback_x,
                y=fallback_y,
                color=valid_color,
                markers=True,
                template=template,
                height=height,
            )
        elif chart_type == "Area Chart":
            fig = px.area(
                working_df,
                x=fallback_x,
                y=fallback_y,
                color=valid_color,
                template=template,
                height=height,
            )
        elif chart_type == "Box Plot":
            fig = px.box(
                working_df,
                x=fallback_x,
                y=fallback_y if fallback_y in all_cols else None,
                color=valid_color,
                template=template,
                height=height,
            )
        elif chart_type == "Violin Plot":
            fig = px.violin(
                working_df,
                x=fallback_x,
                y=fallback_y if fallback_y in all_cols else None,
                color=valid_color,
                box=True,
                points="all",
                template=template,
                height=height,
            )
        elif chart_type == "Pie / Donut":
            counts = working_df[fallback_x].value_counts().reset_index()
            counts.columns = [fallback_x, "count"]
            fig = px.pie(
                counts,
                names=fallback_x,
                values="count",
                hole=0.4,
                template=template,
                height=height,
            )
        elif chart_type == "Heatmap":
            num_df = working_df.select_dtypes(include=[np.number])
            if not num_df.empty and num_df.shape[1] > 1:
                corr = num_df.corr()
                fig = px.imshow(
                    corr,
                    text_auto=True,
                    color_continuous_scale="Viridis",
                    template=template,
                    height=height,
                )
            else:
                fig = px.bar(working_df, template=template, height=height)
        elif chart_type == "Treemap":
            vals = (
                y
                if y and y in num_cols
                else (fallback_y if fallback_y in num_cols else None)
            )
            fig = px.treemap(
                working_df,
                path=[fallback_x],
                values=vals,
                template=template,
                height=height,
            )
        elif chart_type == "Sunburst":
            vals = (
                y
                if y and y in num_cols
                else (fallback_y if fallback_y in num_cols else None)
            )
            path = (
                [fallback_x, valid_color]
                if valid_color and valid_color != fallback_x
                else [fallback_x]
            )
            fig = px.sunburst(
                working_df,
                path=path,
                values=vals,
                template=template,
                height=height,
            )
        elif chart_type == "Funnel Chart":
            fig = px.funnel(
                working_df,
                x=fallback_x,
                y=fallback_y,
                color=valid_color,
                template=template,
                height=height,
            )
        else:
            fig = px.bar(working_df, template=template, height=height)

        fig.update_layout(
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            font=dict(color="#f8fafc", family="Inter, sans-serif"),
            margin=dict(l=20, r=20, t=40, b=20),
            legend=dict(
                orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1
            ),
        )
        return fig
    except Exception as e:
        st.error(f"⚠️ Chart execution error: {e}")
        return None


def render_custom_builder(df: pd.DataFrame):
    section_header(
        "🎨 Advanced Custom Chart Studio",
        "Configure multi-dimensional data visualizations with real-time Plotly rendering & direct exports.",
    )

    if not PLOTLY_AVAILABLE:
        st.error("⚠️ Plotly is required for visualization rendering.")
        return

    all_cols = df.columns.tolist()
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()

    chart_type = st.selectbox(
        "Select Visualization Archetype",
        [
            "Bar Chart",
            "Scatter Plot",
            "Line Chart",
            "Area Chart",
            "Histogram",
            "Box Plot",
            "Violin Plot",
            "Pie / Donut",
            "Heatmap",
            "Treemap",
            "Sunburst",
            "Funnel Chart",
        ],
        key="viz_type_advanced",
    )

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        x_col = st.selectbox("X-Axis / Dimension", [""] + all_cols, key="viz_x_adv")
    with col2:
        y_col = st.selectbox("Y-Axis / Metric", [""] + all_cols, key="viz_y_adv")
    with col3:
        color_col = st.selectbox(
            "Color / Group By", [""] + all_cols, key="viz_color_adv"
        )
    with col4:
        size_col = st.selectbox(
            "Marker Size (Scatter)", [""] + numeric_cols, key="viz_size_adv"
        )

    filter_outliers = st.checkbox(
        "🧹 Automatically Filter Extreme Outliers (3 Sigma Rule on Numeric Metrics)",
        value=False,
        key="viz_outlier_filter",
    )

    render_df = df.copy()
    if filter_outliers and y_col and y_col in numeric_cols:
        series = render_df[y_col].dropna()
        mean, std = series.mean(), series.std()
        if std > 0:
            render_df = render_df[(render_df[y_col] - mean).abs() <= 3 * std]

    height = st.slider(
        "Visualization Height (px)",
        300,
        750,
        450,
        25,
        key="viz_height_adv",
    )

    fig = build_chart(
        chart_type,
        render_df,
        x=x_col if x_col else None,
        y=y_col if y_col else None,
        color=color_col if color_col else None,
        size=size_col if size_col else None,
        height=height,
    )

    if fig:
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": True})

        if y_col and y_col in numeric_cols:
            st.info(
                f"📈 **Active Metric Telemetry (`{y_col}`)**: Count: `{len(render_df)}` | "
                f"Mean: `{render_df[y_col].mean():,.2f}` | Median: `{render_df[y_col].median():,.2f}` | "
                f"Std Dev: `{render_df[y_col].std():,.2f}`"
            )

        st.markdown("#### 📥 Download & Copy Studio Assets")
        exp_col1, exp_col2 = st.columns(2)
        with exp_col1:
            render_export_buttons(
                render_df,
                base_name=f"custom_viz_{chart_type.lower().replace(' ', '_')}",
            )
        with exp_col2:
            csv_data = render_df.to_csv(index=False).encode("utf-8")
            st.download_button(
                label="📋 Copy / Download Raw Table CSV",
                data=csv_data,
                file_name="chart_source_data.csv",
                mime="text/csv",
                key="download_raw_csv_custom",
            )
    else:
        st.error("⚠️ Could not render chart with the selected parameters.")


def _skewness(series: pd.Series) -> float:
    """Compute Fisher-Pearson coefficient of skewness safely."""
    s = series.dropna()
    if len(s) < 3 or s.std() == 0:
        return 0.0
    return float(((s - s.mean()) ** 3).mean() / (s.std() ** 3))


def render_auto_studio(df: pd.DataFrame):
    section_header(
        "🤖 AI Auto-Recommendation Studio",
        "Automated exploratory visual discovery — selections are driven by actual data properties (correlation strength, skew, cardinality, temporal structure).",
    )

    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()
    datetime_cols = df.select_dtypes(
        include=["datetime64[ns]", "datetime64[ns, UTC]"]
    ).columns.tolist()

    # Cast string dates dynamically if not parsed yet
    if not datetime_cols:
        for c in cat_cols:
            if "date" in c.lower() or "time" in c.lower():
                try:
                    df[c] = pd.to_datetime(df[c])
                    datetime_cols.append(c)
                    break
                except (ValueError, TypeError):
                    pass

    low_card_cats = [c for c in cat_cols if df[c].nunique() <= 15]

    if not numeric_cols:
        st.warning("⚠️ Dataset requires numeric columns to generate automated visual recommendations.")
        return

    recs = []
    for col in numeric_cols[:2]:
        skew = _skewness(df[col])
        if abs(skew) > 1.0:
            recs.append(
                ("Box Plot", {"y": col}, f"Distribution of '{col}' — skew = {skew:.2f}, box plot recommended.")
            )
        else:
            recs.append(
                ("Histogram", {"x": col}, f"Univariate Distribution Analysis of '{col}' (skew = {skew:.2f})")
            )

    if len(numeric_cols) >= 2:
        corr = df[numeric_cols].corr().abs()
        upper = corr.where(np.triu(np.ones(corr.shape), k=1).astype(bool))
        stacked = upper.stack()
        if not stacked.empty:
            (best_x, best_y), best_r = stacked.idxmax(), stacked.max()
            recs.append(
                (
                    "Scatter Plot",
                    {"x": best_x, "y": best_y},
                    f"Strongest Bivariate Relationship: '{best_x}' vs '{best_y}' (|r| = {best_r:.2f}).",
                )
            )

    if datetime_cols and numeric_cols:
        recs.append(
            (
                "Line Chart",
                {"x": datetime_cols[0], "y": numeric_cols[0]},
                f"Temporal Trend of '{numeric_cols[0]}' over '{datetime_cols[0]}'.",
            )
        )

    if low_card_cats and numeric_cols:
        best_cat = min(low_card_cats, key=lambda c: df[c].nunique())
        recs.append(
            (
                "Bar Chart",
                {"x": best_cat, "y": numeric_cols[0]},
                f"Categorical Comparison of '{numeric_cols[0]}' across '{best_cat}'.",
            )
        )

    for i, (ctype, params, rationale) in enumerate(recs):
        with st.container():
            st.markdown(f"#### 💡 Insight Perspective {i+1}: {ctype}")
            st.caption(rationale)
            fig = build_chart(ctype, df, height=350, **params)
            if fig:
                st.plotly_chart(fig, use_container_width=True)
            st.markdown("---")


def render_exec_dashboard(df: pd.DataFrame):
    section_header(
        "📊 Executive KPI & Multi-Chart Dashboard",
        "Assembled executive dashboard combining core metrics, distribution telemetry, and trend analytics.",
    )

    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()

    if not numeric_cols:
        st.warning("⚠️ Insufficient numeric features for executive metrics.")
        return

    st.markdown("#### Key Performance Indicators (KPIs)")
    kpi_cols = st.columns(min(4, len(numeric_cols)))
    for i, col_name in enumerate(numeric_cols[:4]):
        mean_val = df[col_name].mean()
        std_val = df[col_name].std()
        kpi_cols[i].metric(
            label=f"Avg {col_name}",
            value=f"{mean_val:,.2f}",
            delta=f"±{std_val:.2f} σ",
        )

    st.markdown("---")
    st.markdown("#### Multi-Panel Executive Telemetry")
    row1_c1, row1_c2 = st.columns(2)

    with row1_c1:
        fig1 = build_chart("Histogram", df, x=numeric_cols[0], height=340)
        if fig1:
            st.plotly_chart(fig1, use_container_width=True)

    with row1_c2:
        if len(numeric_cols) >= 2:
            fig2 = build_chart(
                "Scatter Plot", df, x=numeric_cols[0], y=numeric_cols[1], height=340
            )
            if fig2:
                st.plotly_chart(fig2, use_container_width=True)

    if cat_cols and len(numeric_cols) >= 1:
        fig3 = build_chart(
            "Bar Chart", df, x=cat_cols[0], y=numeric_cols[0], height=340
        )
        if fig3:
            st.plotly_chart(fig3, use_container_width=True)

    st.markdown("---")
    render_export_buttons(df, base_name="executive_dashboard_dataset")


def _fig_to_png_bytes(fig: go.Figure) -> bytes:
    """Render Plotly figure directly to PNG bytes."""
    return fig.to_image(format="png", width=1000, height=560, scale=2)


def _build_pptx(deck_title: str, slides_spec: list) -> bytes:
    """Generates standard .pptx binary presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = deck_title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = (
            f"Generated {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}"
        )

    blank_layout = prs.slide_layouts[6]
    for spec in slides_spec:
        slide = prs.slides.add_slide(blank_layout)
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
        )
        tb.text_frame.text = spec["title"]
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        tb.text_frame.paragraphs[0].font.bold = True

        metric_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.6)
        )
        metric_box.text_frame.text = (
            f"{spec['metric_label']}: {spec['metric_value']}"
        )
        metric_box.text_frame.paragraphs[0].font.size = Pt(18)

        if spec.get("image_bytes"):
            slide.shapes.add_picture(
                io.BytesIO(spec["image_bytes"]),
                Inches(0.7),
                Inches(1.9),
                width=Inches(11.9),
            )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_deck_builder(df: pd.DataFrame):
    section_header(
        "📹 Presentation Deck & Slide Generator",
        "Structure executive presentation slides with embedded charts and narrative metric highlights — with downloadable .pptx output.",
    )

    deck_title = st.text_input(
        "Presentation Deck Title",
        value="Executive Analytics & Data Briefing",
        key="deck_title_input",
    )
    slide_count = st.slider(
        "Number of Slides to Generate", 2, 6, 4, key="deck_slide_count"
    )

    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()

    if not PPTX_AVAILABLE:
        st.warning(
            "⚠️ `.pptx` export library missing. Visual previews are available on-screen."
        )
    if not KALEIDO_AVAILABLE:
        st.info(
            "ℹ️ Image rendering engine missing (`kaleido`). Slides will export metrics with text layout until available."
        )

    if st.button(
        "📹 Generate Presentation Deck", type="primary", key="build_deck_btn"
    ):
        if not numeric_cols:
            st.error("Need at least one numeric column to build slide metrics.")
            return

        slides_spec = []
        for i in range(slide_count):
            col_metric = numeric_cols[i % len(numeric_cols)]
            mean_val = df[col_metric].mean()
            image_bytes = None
            fig = None
            if cat_cols:
                fig = build_chart(
                    "Bar Chart",
                    df,
                    x=cat_cols[i % len(cat_cols)],
                    y=col_metric,
                    height=350,
                )
            if fig and KALEIDO_AVAILABLE:
                try:
                    image_bytes = _fig_to_png_bytes(fig)
                except Exception as e:
                    st.caption(f"Slide {i+1}: chart image render skipped ({e})")

            slides_spec.append(
                {
                    "title": f"Slide {i+1}: {col_metric} Briefing",
                    "metric_label": f"Average {col_metric}",
                    "metric_value": f"{mean_val:,.2f}",
                    "image_bytes": image_bytes,
                    "fig": fig,
                }
            )

        st.markdown("#### On-Screen Preview")
        for i, spec in enumerate(slides_spec):
            with st.expander(f"{spec['title']}", expanded=(i == 0)):
                st.metric(spec["metric_label"], spec["metric_value"])
                if spec["fig"]:
                    st.plotly_chart(spec["fig"], use_container_width=True)

        if PPTX_AVAILABLE:
            pptx_bytes = _build_pptx(deck_title, slides_spec)
            st.success(
                f"✅ Presentation deck '{deck_title}' compiled — {slide_count} slides ready."
            )
            st.download_button(
                "⬇️ Download Presentation (.pptx)",
                data=pptx_bytes,
                file_name=f"{deck_title.lower().replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="dl_pptx_deck",
            )


def render_chart_extractor(df: pd.DataFrame):
    section_header(
        "📊 Chart Data Extractor & Aggregator",
        "Extract, aggregate, and export structured tabular subsets directly from visual nodes.",
    )

    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()

    if not numeric_cols or not cat_cols:
        st.warning("⚠️ Data extraction requires both categorical and numeric columns.")
        return

    col1, col2, col3 = st.columns(3)
    with col1:
        group_col = st.selectbox("Group By Dimension", cat_cols, key="ext_group")
    with col2:
        agg_col = st.selectbox("Target Metric", numeric_cols, key="ext_agg")
    with col3:
        agg_func = st.selectbox(
            "Aggregation Function",
            ["mean", "sum", "median", "count", "max", "min", "std"],
            key="ext_func",
        )

    if st.button(
        "📊 Extract & Process Aggregated Data", type="primary", key="run_extraction"
    ):
        if agg_func == "count":
            extracted_df = df.groupby(group_col)[agg_col].count().reset_index()
        else:
            extracted_df = (
                df.groupby(group_col)[agg_col].agg(agg_func).reset_index()
            )

        st.markdown("#### 📋 Extracted Dataset Preview & Copy Options")
        st.dataframe(extracted_df, use_container_width=True, hide_index=True)
        render_export_buttons(extracted_df, base_name="extracted_chart_dataset")


def main():
    from modules.subscription import require_active_subscription

    require_active_subscription(hub_id="visualization")

    setup_page("Visualization Studio", "📈", initial_sidebar_state="expanded")

    from modules.user_preferences import (
        render_accent_color_css,
        render_readability_fix,
    )

    render_readability_fix()
    render_accent_color_css()

    hero_card(
        "📈 Visualization Studio & Executive Dashboard Hub (Premium)",
        "Consolidated enterprise visualization platform featuring advanced multi-dimensional charting, a data-driven AI recommendation engine, executive dashboards, a downloadable presentation deck generator, and chart data extraction.",
        badge_text="VISUALIZATION STUDIO • PREMIUM TIER",
    )

    render_dataset_context_banner()
    df = get_df()

    tabs = st.tabs(
        [
            "🎨 Custom Builder",
            "🤖 Auto-Recommendations",
            "📊 Executive Dashboard",
            "📹 Presentation Deck",
            "📊 Chart Data Extractor",
        ]
    )

    with tabs[0]:
        render_custom_builder(df)
    with tabs[1]:
        render_auto_studio(df)
    with tabs[2]:
        render_exec_dashboard(df)
    with tabs[3]:
        render_deck_builder(df)
    with tabs[4]:
        render_chart_extractor(df)

    render_standard_footer("VISUALIZATION STUDIO")


if __name__ == "__main__":
    main()
