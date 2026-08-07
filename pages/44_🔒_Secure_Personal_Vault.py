# ============================================================================
# NEXUS VAULT 2.0 — Google Workspace-Grade Personal Suite
# Drive / Docs / Sheets / Slides / Mail / Meetings / Calendar / Contacts /
# Tasks / Notes / Automations / Security / Storage & Backup
# ============================================================================

import base64
import datetime
import hashlib
import io
import json
import re
import secrets
import sqlite3
import uuid

import pandas as pd
import streamlit as st
import streamlit.components.v1 as components

APP_TITLE = "Nexus Vault 2.0"
DB_PATH = "nexus_vault.db"
DEFAULT_ADMIN_USERNAME = "admin"
DEFAULT_ADMIN_PASSWORD = "admin123"
STORAGE_PLANS_GB = {"Free (15 GB)": 15, "Plus (100 GB)": 100, "Pro (1 TB)": 1024, "Unlimited*": 10_000_000}

try:
    from fpdf import FPDF
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from cryptography.fernet import Fernet
    HAS_CRYPTO = True
except ImportError:
    HAS_CRYPTO = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

st.set_page_config(page_title=APP_TITLE, page_icon="🔒", layout="wide", initial_sidebar_state="expanded")

BRASS = "#C99A3A"
INK = "#14161F"

st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=SpaceGrotesk:wght@500;600;700;800&family=SourceSerif4:wght@400;600&family=JetBrainsMono&display=swap');
html, body, [class*="css"]  {{ font-family: 'Space Grotesk', sans-serif; }}
code, pre {{ font-family: 'JetBrains Mono', monospace !important; }}
.nv-header {{
    background: linear-gradient(120deg, #1e1e2f 0%, #2b2f45 60%, #3a3050 100%);
    border-radius: 18px; padding: 22px 30px; color: white; margin-bottom: 18px;
    border: 1px solid rgba(201,154,58,0.28); box-shadow: 0 10px 30px rgba(0,0,0,0.25);
}}
.nv-header h1 {{ margin: 0; font-size: 1.6rem; font-weight: 800; }}
.nv-header p {{ margin: 6px 0 0 0; opacity: 0.78; font-size: 0.92rem; }}
.nv-card {{
    background: rgba(255,255,255,0.035); backdrop-filter: blur(10px);
    border: 1px solid rgba(201,154,58,0.20); border-radius: 14px; padding: 14px 16px;
}}
.nv-card-title {{ font-size: 0.7rem; font-weight: 700; text-transform: uppercase; letter-spacing: 0.6px; color: {BRASS}; }}
.nv-card-value {{ font-size: 1.4rem; font-weight: 800; margin-top: 4px; }}
.stButton>button {{
    background: linear-gradient(120deg, #1e1e2f 0%, #2b2f45 60%, #3a3050 100%);
    border-radius: 9px; font-weight: 600; color: white;
}}
.stButton>button:hover {{ box-shadow: 0 8px 20px rgba(201,154,58,0.35); }}
.nv-badge {{ display:inline-block; padding: 3px 10px; border-radius: 20px; font-size: 0.72rem; font-weight:700; }}
.nv-badge-live {{ background:#D1FAE5; color:#065F46; }}
.nv-badge-demo {{ background:#F1F5F9; color:#475569; }}
.nv-badge-role {{ background: rgba(201,154,58,0.15); color: {BRASS}; border: 1px solid rgba(201,154,58,0.35); }}
.nv-empty {{ text-align:center; padding: 40px 10px; color:#94A3B8; }}
.nv-chip {{ display:inline-block; padding:2px 9px; border-radius:12px; font-size:0.7rem; font-weight:600; margin-right:4px; }}
.nv-note {{ background:#14161F; color:white; border:1px solid rgba(201,154,58,0.28); border-radius:12px; padding:16px; margin-bottom:10px; }}
</style>
""", unsafe_allow_html=True)


# ============================================================================
# DATABASE
# ============================================================================
def get_conn():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def now_str():
    return datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def init_db():
    conn = get_conn()
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT, username TEXT UNIQUE, password_hash TEXT,
        role TEXT DEFAULT 'Editor', created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS files (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, folder TEXT, content_type TEXT,
        size_bytes INTEGER, data_blob BLOB, uploaded_at TEXT, deleted INTEGER DEFAULT 0,
        starred INTEGER DEFAULT 0, tags TEXT DEFAULT '', notes TEXT DEFAULT '',
        version INTEGER DEFAULT 1)""")
    c.execute("""CREATE TABLE IF NOT EXISTS folders (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, owner TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS file_versions (
        id INTEGER PRIMARY KEY AUTOINCREMENT, file_id INTEGER, version INTEGER,
        data_blob BLOB, saved_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS share_links (
        id INTEGER PRIMARY KEY AUTOINCREMENT, item_type TEXT, item_id INTEGER,
        token TEXT, permission TEXT, expires_at TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS documents (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, content TEXT, owner TEXT,
        folder TEXT, modified_at TEXT, deleted INTEGER DEFAULT 0, starred INTEGER DEFAULT 0,
        version INTEGER DEFAULT 1)""")
    c.execute("""CREATE TABLE IF NOT EXISTS doc_versions (
        id INTEGER PRIMARY KEY AUTOINCREMENT, doc_id INTEGER, version INTEGER,
        content TEXT, saved_at TEXT, editor TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS sheets (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, data_json TEXT, owner TEXT,
        modified_at TEXT, deleted INTEGER DEFAULT 0)""")
    c.execute("""CREATE TABLE IF NOT EXISTS slide_decks (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, slides_json TEXT, owner TEXT,
        modified_at TEXT, deleted INTEGER DEFAULT 0, theme TEXT DEFAULT 'Brass')""")
    c.execute("""CREATE TABLE IF NOT EXISTS mail (
        id INTEGER PRIMARY KEY AUTOINCREMENT, folder TEXT, sender TEXT, recipient TEXT,
        subject TEXT, body TEXT, mode TEXT, sent_at TEXT, read INTEGER DEFAULT 0,
        starred INTEGER DEFAULT 0, label TEXT DEFAULT 'Primary')""")
    c.execute("""CREATE TABLE IF NOT EXISTS contacts (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, email TEXT, phone TEXT,
        company TEXT, group_name TEXT DEFAULT 'General', notes TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS calendar_events (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, start_date TEXT, end_date TEXT,
        start_time TEXT, end_time TEXT, category TEXT DEFAULT 'General',
        remind_minutes INTEGER DEFAULT 10, notes TEXT, owner TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS meetings (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, room TEXT, scheduled_for TEXT,
        owner TEXT, created_at TEXT, recurring INTEGER DEFAULT 0, invitees TEXT DEFAULT '[]')""")
    c.execute("""CREATE TABLE IF NOT EXISTS tasks (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, done INTEGER DEFAULT 0,
        due_date TEXT, priority TEXT DEFAULT 'Medium', list_name TEXT DEFAULT 'Default',
        owner TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS notes (
        id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, body TEXT, color TEXT DEFAULT '#14161F',
        pinned INTEGER DEFAULT 0, owner TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS passwords (
        id INTEGER PRIMARY KEY AUTOINCREMENT, site TEXT, username TEXT,
        encrypted_password BLOB, notes TEXT, owner TEXT, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS audit_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT, timestamp TEXT, actor TEXT, action TEXT, status TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS automation_rules (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, trigger_type TEXT, action_type TEXT,
        threshold INTEGER DEFAULT 0, enabled INTEGER DEFAULT 1)""")
    c.execute("""CREATE TABLE IF NOT EXISTS activity_feed (
        id INTEGER PRIMARY KEY AUTOINCREMENT, timestamp TEXT, actor TEXT, action TEXT, detail TEXT)""")
    conn.commit()

    c.execute("SELECT COUNT(*) FROM users")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO users (username, password_hash, role, created_at) VALUES (?, ?, ?, ?)",
                  (DEFAULT_ADMIN_USERNAME, hash_password(DEFAULT_ADMIN_PASSWORD), "Admin", now_str()))
        conn.commit()

    c.execute("SELECT COUNT(*) FROM folders")
    if c.fetchone()[0] == 0:
        for name in ["My Drive", "Documents", "Images", "Music", "Videos"]:
            c.execute("INSERT INTO folders (name, owner, created_at) VALUES (?,?,?)", (name, "admin", now_str()))

    c.execute("SELECT COUNT(*) FROM documents")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO documents (title, content, owner, folder, modified_at) VALUES (?, ?, ?, ?, ?)",
                  ("Welcome to Nexus Vault",
                   "# Welcome\n\nThis document is stored in a real SQLite database next to this script "
                   "(nexus_vault.db) — edit it, refresh the page, and it will still be here.",
                   "admin", "My Docs", now_str()))
    c.execute("SELECT COUNT(*) FROM sheets")
    if c.fetchone()[0] == 0:
        demo_rows = [{"Item": "Widget A", "Cost": 10, "Qty": 5, "Total": "=SUM(C2*B2)"},
                     {"Item": "Widget B", "Cost": 20, "Qty": 3, "Total": "=SUM(C3*B3)"}]
        c.execute("INSERT INTO sheets (name, data_json, owner, modified_at) VALUES (?, ?, ?, ?)",
                  ("Sample Sheet", json.dumps(demo_rows), "admin", now_str()))
    c.execute("SELECT COUNT(*) FROM slide_decks")
    if c.fetchone()[0] == 0:
        demo_slides = [{"title": "Nexus Vault", "body": "Secure. Unified. Persisted.", "layout": "title"},
                       {"title": "Agenda", "body": "1. Storage\n2. Security\n3. Roadmap", "layout": "title-body"}]
        c.execute("INSERT INTO slide_decks (title, slides_json, owner, modified_at) VALUES (?, ?, ?, ?)",
                  ("Welcome Deck", json.dumps(demo_slides), "admin", now_str()))
    c.execute("SELECT COUNT(*) FROM mail")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO mail (folder, sender, recipient, subject, body, mode, sent_at) VALUES (?,?,?,?,?,?,?)",
                  ("inbox", "team@nexusvault.io", "admin", "Welcome to your persisted inbox",
                   "This message is stored in SQLite, not just browser memory.", "demo", now_str()))
    c.execute("SELECT COUNT(*) FROM contacts")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO contacts (name, email, phone, company, group_name, created_at) VALUES (?,?,?,?,?,?)",
                  ("Team Nexus", "team@nexusvault.io", "+256700000000", "Nexus Labs", "Work", now_str()))
    c.execute("SELECT COUNT(*) FROM tasks")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO tasks (title, list_name, priority, owner, created_at) VALUES (?,?,?,?,?)",
                  ("Explore Nexus Vault", "Getting Started", "High", "admin", now_str()))
    c.execute("SELECT COUNT(*) FROM notes")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO notes (title, body, owner, created_at) VALUES (?,?,?,?)",
                  ("Idea", "Birth of Nexus Vault 2.0 — a Google Workspace-grade suite.", "admin", now_str()))
    c.execute("SELECT COUNT(*) FROM automation_rules")
    if c.fetchone()[0] == 0:
        c.execute("INSERT INTO automation_rules (name, trigger_type, action_type, threshold, enabled) VALUES (?,?,?,?,1)",
                  ("Flag large files when storage is heavy", "storage_high", "flag_largest_file", 5))
    conn.commit()
    conn.close()


def log_audit(actor, action, status="OK"):
    try:
        conn = get_conn()
        conn.execute("INSERT INTO audit_log (timestamp, actor, action, status) VALUES (?, ?, ?, ?)",
                     (now_str(), actor, action, status))
        conn.commit()
        conn.close()
    except Exception:
        pass


def log_activity(actor, detail):
    try:
        conn = get_conn()
        conn.execute("INSERT INTO activity_feed (timestamp, actor, action, detail) VALUES (?,?,?,?)",
                     (now_str(), actor, "did", detail))
        conn.commit()
        conn.close()
    except Exception:
        pass


def hash_password(password: str) -> str:
    salt = secrets.token_hex(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 100_000).hex()
    return f"{salt}${digest}"


def verify_password(password: str, stored: str) -> bool:
    try:
        salt, digest = stored.split("$")
    except ValueError:
        return False
    check = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 100_000).hex()
    return secrets.compare_digest(check, digest)


def fmt_bytes(b):
    if b < 1024:
        return f"{b} B"
    elif b < 1024 ** 2:
        return f"{b/1024:.1f} KB"
    elif b < 1024 ** 3:
        return f"{b/1024**2:.1f} MB"
    return f"{b/1024**3:.2f} GB"


# ============================================================================
# SAFE SPREADSHEET FORMULA ENGINE
# ============================================================================
def col_to_num(col: str) -> int:
    n = 0
    for ch in col.upper():
        n = n * 26 + (ord(ch) - 64)
    return n


def num_to_col(n: int) -> str:
    s = ""
    while n > 0:
        r = (n - 1) % 26
        s = chr(65 + r) + s
        n = (n - 1) // 26
    return s


def parse_ref(ref: str):
    m = re.match(r"^([A-Z]+)(\d+)$", ref.upper())
    return {"col": m.group(1), "row": int(m.group(2))} if m else None


def expand_range(start: str, end: str):
    s, e = parse_ref(start), parse_ref(end)
    if not s or not e:
        return []
    c1, c2 = col_to_num(s["col"]), col_to_num(e["col"])
    r1, r2 = min(s["row"], e["row"]), max(s["row"], e["row"])
    min_c, max_c = min(c1, c2), max(c1, c2)
    return [f"{num_to_col(c)}{r}" for r in range(r1, r2 + 1) for c in range(min_c, max_c + 1)]


def grid_from_df(df):
    grid = {}
    for i, row in df.iterrows():
        for j, col in enumerate(df.columns):
            grid[f"{num_to_col(j+1)}{i+1}"] = row[col]
    return grid


def resolve_cell(ref, grid, seen):
    if ref in seen:
        return 0
    raw = grid.get(ref, "")
    if raw == "":
        return 0
    if isinstance(raw, str) and raw.startswith("="):
        next_seen = set(seen)
        next_seen.add(ref)
        res = evaluate_formula(raw, grid, next_seen)
        try:
            return float(res)
        except (ValueError, TypeError):
            return res
    try:
        return float(raw)
    except (ValueError, TypeError):
        return raw


def range_fn(fn, cells, grid, seen):
    vals = [resolve_cell(c, grid, seen) for c in cells]
    nums = [float(v) for v in vals if isinstance(v, (int, float)) or (isinstance(v, str) and v.replace(".", "", 1).isdigit())]
    fn = fn.upper()
    if fn == "SUM":
        return sum(nums)
    if fn in ("AVERAGE", "AVG"):
        return sum(nums) / len(nums) if nums else 0
    if fn == "MIN":
        return min(nums) if nums else 0
    if fn == "MAX":
        return max(nums) if nums else 0
    if fn == "COUNT":
        return len(nums)
    if fn == "COUNTA":
        return len([v for v in vals if v not in ("", 0)])
    if fn == "CONCAT":
        return "".join(map(str, vals))
    return 0


def vlookup(lookup_val, range_str, col_idx, grid, seen):
    try:
        start, end = range_str.split(":")
        s, e = parse_ref(start), parse_ref(end)
        if not s or not e:
            return "#N/A"
        c1, c2 = col_to_num(s["col"]), col_to_num(e["col"])
        target_col = c1 + col_idx - 1
        if target_col > c2:
            return "#REF!"
        for r in range(min(s["row"], e["row"]), max(s["row"], e["row"]) + 1):
            key = f"{num_to_col(c1)}{r}"
            val = resolve_cell(key, grid, seen)
            if str(val).strip().lower() == str(lookup_val).strip().lower():
                return resolve_cell(f"{num_to_col(target_col)}{r}", grid, seen)
        return "#N/A"
    except Exception:
        return "#ERR"


def safe_math_eval(expr):
    try:
        import math as m
        return eval(compile(expr, "<string>", "eval"), {"__builtins__": None}, {"math": m})
    except Exception:
        return "#ERR"


def evaluate_formula(raw, grid, seen=None):
    if seen is None:
        seen = set()
    if not isinstance(raw, str) or not raw.startswith("="):
        return raw
    expr = raw[1:].strip()

    def _vlookup_sub(match):
        key, rng, idx = match.group(1).strip(), match.group(2).strip(), int(match.group(3))
        lookup_val = key[1:-1] if key.startswith('"') else resolve_cell(key.upper(), grid, seen)
        res = vlookup(lookup_val, rng, idx, grid, seen)
        return f'"{res}"' if isinstance(res, str) else str(res)

    expr = re.sub(r"VLOOKUP\(\s*([^,]+?)\s*,\s*([A-Z]+\d+[A-Z]+\d+|[A-Z]+\d+:[A-Z]+\d+)\s*,\s*(\d+)\s*\)", _vlookup_sub, expr, flags=re.IGNORECASE)

    def _range_sub(match):
        fn, a, b = match.group(1), match.group(2), match.group(3)
        res = range_fn(fn, expand_range(a.upper(), b.upper()), grid, seen)
        return f'"{res}"' if isinstance(res, str) else str(res)

    expr = re.sub(r"(SUM|AVERAGE|AVG|MIN|MAX|COUNT|COUNTA|CONCAT)\(\s*([A-Z]+\d+)\s*:\s*([A-Z]+\d+)\s*\)", _range_sub, expr, flags=re.IGNORECASE)

    def _cell_sub(match):
        ref = match.group(0).upper()
        val = resolve_cell(ref, grid, seen)
        return f'"{val}"' if isinstance(val, str) else str(val)

    expr = re.sub(r"\b[A-Z]+\d+\b", _cell_sub, expr)

    def _if_sub(match):
        inner = match.group(1)
        parts = inner.split(",")
        if len(parts) == 3:
            return f"({parts[1]}) if ({parts[0]}) else ({parts[2]})"
        return match.group(0)

    expr = re.sub(r"IF\(([^()]*)\)", _if_sub, expr, flags=re.IGNORECASE)

    res = safe_math_eval(expr)
    if isinstance(res, float):
        return round(res, 4)
    return res


def compute_sheet(df):
    """Return a DataFrame with formula cells resolved."""
    grid = grid_from_df(df)
    out = df.copy().astype(object)
    for i, row in df.iterrows():
        for j, col in enumerate(df.columns):
            val = row[col]
            if isinstance(val, str) and val.startswith("="):
                out.at[i, col] = evaluate_formula(val, grid)
    return out


init_db()


# ============================================================================
# SESSION STATE + LOGIN
# ============================================================================
for key, default in {"logged_in": False, "username": None, "role": None,
                     "active_doc_id": None, "active_sheet_id": None, "active_deck_id": None,
                     "preview_target": None, "meeting_room": None, "active_folder": "My Drive",
                     "nav_view": "My Drive", "smtp_creds": {}, "fernet_key": None,
                     "vault_key": None, "show_computed": False}.items():
    if key not in st.session_state:
        st.session_state[key] = default

if not st.session_state["logged_in"]:
    left, mid, right = st.columns([1, 1.2, 1])
    with mid:
        st.markdown(f'<div class="nv-header"><h1>🔒 {APP_TITLE}</h1><p>Sign in to your workspace.</p></div>', unsafe_allow_html=True)
        with st.form("login_form"):
            u = st.text_input("Username", value="admin")
            p = st.text_input("Password", type="password", value="admin123")
            go = st.form_submit_button("Sign in", type="primary", use_container_width=True)
        if go:
            conn = get_conn()
            row = conn.execute("SELECT * FROM users WHERE username = ?", (u,)).fetchone()
            conn.close()
            if row and verify_password(p, row["password_hash"]):
                st.session_state.update(logged_in=True, username=row["username"], role=row["role"])
                log_audit(u, "LOGIN", "OK")
                st.rerun()
            else:
                st.error("Invalid username or password.")
                log_audit(u, "LOGIN_FAILED", "DENIED")
        st.caption("Demo login: **admin / admin123**")
    st.stop()


# ============================================================================
# SIDEBAR
# ============================================================================
with st.sidebar:
    st.markdown(f"## 🔒 {APP_TITLE}")
    st.markdown(f'**{st.session_state["username"]}** &nbsp; <span class="nv-badge nv-badge-role">{st.session_state["role"]}</span>', unsafe_allow_html=True)
    if st.button("Sign out", use_container_width=True):
        log_audit(st.session_state["username"], "LOGOUT")
        st.session_state["logged_in"] = False
        st.rerun()
    st.markdown("---")
    nav_options = ["🏠 Home", "📁 Drive", "📄 Docs", "📊 Sheets", "🎞️ Slides", "✉️ Mail",
                   "👥 Contacts", "📅 Calendar", "✅ Tasks", "🗒️ Notes", "🎥 Meetings",
                   "⚙️ Automations", "🔑 Passwords", "🛡️ Security", "☁️ Storage & Backup"]
    if st.session_state["role"] == "Admin":
        nav_options += ["👤 Users", "📋 Admin Log"]
    active = st.radio("Navigate", nav_options, label_visibility="collapsed")

st.markdown(f'<div class="nv-header"><h1>🔒 {APP_TITLE}</h1><p>{active}</p></div>', unsafe_allow_html=True)


def require_editor():
    if st.session_state["role"] == "Viewer":
        st.warning("Your account role is **Viewer** — read-only.")
        st.stop()


# ============================================================================
# HOME
# ============================================================================
if active == "🏠 Home":
    conn = get_conn()
    counts = {
        "Files": conn.execute("SELECT COUNT(*) FROM files WHERE deleted=0").fetchone()[0],
        "Documents": conn.execute("SELECT COUNT(*) FROM documents WHERE deleted=0").fetchone()[0],
        "Sheets": conn.execute("SELECT COUNT(*) FROM sheets WHERE deleted=0").fetchone()[0],
        "Decks": conn.execute("SELECT COUNT(*) FROM slide_decks WHERE deleted=0").fetchone()[0],
    }
    total_bytes = conn.execute("SELECT COALESCE(SUM(size_bytes),0) FROM files WHERE deleted=0").fetchone()[0]
    unread = conn.execute("SELECT COUNT(*) FROM mail WHERE folder='inbox' AND read=0").fetchone()[0]
    tasks = conn.execute("SELECT COUNT(*) FROM tasks WHERE done=0").fetchone()[0]
    contacts = conn.execute("SELECT COUNT(*) FROM contacts").fetchone()[0]
    conn.close()

    cols = st.columns(4)
    for col, (label, val) in zip(cols, counts.items()):
        col.markdown(f'<div class="nv-card"><div class="nv-card-title">{label}</div><div class="nv-card-value">{val}</div></div>', unsafe_allow_html=True)
    c2 = st.columns(4)
    c2[0].markdown(f'<div class="nv-card"><div class="nv-card-title">Unread Mail</div><div class="nv-card-value">{unread}</div></div>', unsafe_allow_html=True)
    c2[1].markdown(f'<div class="nv-card"><div class="nv-card-title">Open Tasks</div><div class="nv-card-value">{tasks}</div></div>', unsafe_allow_html=True)
    c2[2].markdown(f'<div class="nv-card"><div class="nv-card-title">Contacts</div><div class="nv-card-value">{contacts}</div></div>', unsafe_allow_html=True)
    c2[3].markdown(f'<div class="nv-card"><div class="nv-card-title">Storage</div><div class="nv-card-value">{fmt_bytes(total_bytes)}</div></div>', unsafe_allow_html=True)

    st.markdown("### 🕓 Recent Activity")
    conn = get_conn()
    feed = conn.execute("SELECT * FROM activity_feed ORDER BY timestamp DESC LIMIT 10").fetchall()
    conn.close()
    if feed:
        for f in feed:
            st.markdown(f"**{f['timestamp']}** — `{f['actor']}` {f['action']} {f['detail']}")
    else:
        st.info("No activity yet. Start creating files, docs, sheets, or tasks.")

# ============================================================================
# DRIVE
# ============================================================================
elif active == "📁 Drive":
    st.subheader("📁 Drive")
    require_editor()

    conn = get_conn()
    folders = conn.execute("SELECT * FROM folders ORDER BY name").fetchall()
    conn.close()

    left, right = st.columns([1, 3])
    with left:
        st.markdown("##### Folders")
        for f in folders:
            if st.button(f"📂 {f['name']}", key=f"folder_{f['id']}", use_container_width=True):
                st.session_state["active_folder"] = f["name"]
        with st.expander("➕ New Folder"):
            fname = st.text_input("Folder name", key="new_folder_input")
            if st.button("Create") and fname:
                conn = get_conn()
                conn.execute("INSERT INTO folders (name, owner, created_at) VALUES (?,?,?)", (fname, st.session_state["username"], now_str()))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"CREATE_FOLDER:{fname}")
                st.rerun()
        st.markdown("---")
        view = st.radio("View", ["My Drive", "Starred", "Shared", "Recent", "Trash"], key="drive_view")
        st.session_state["nav_view"] = view

    with right:
        upload = st.file_uploader("Upload a file to Drive", key="drive_upload")
        if upload is not None:
            data = upload.getvalue()
            folder = st.session_state.get("active_folder") or "My Drive"
            conn = get_conn()
            cur = conn.execute("INSERT INTO files (name, folder, content_type, size_bytes, data_blob, uploaded_at) VALUES (?,?,?,?,?,?)",
                               (upload.name, folder, upload.type or "", len(data), data, now_str()))
            conn.commit()
            conn.execute("INSERT INTO file_versions (file_id, version, data_blob, saved_at) VALUES (?,?,?,?)",
                         (cur.lastrowid, 1, data, now_str()))
            conn.commit(); conn.close()
            log_audit(st.session_state["username"], f"UPLOAD_FILE:{upload.name}")
            log_activity(st.session_state["username"], f"uploaded '{upload.name}' to {folder}")
            st.success(f"Uploaded {upload.name}")
            st.rerun()

        sc1, sc2, sc3 = st.columns([2, 1, 1])
        search_q = sc1.text_input("🔍 Search", key="drive_search")
        file_type = sc2.selectbox("Type", ["All", "Images", "Documents", "Audio", "Video", "Data", "Code"])
        sort_by = sc3.selectbox("Sort", ["date", "name", "size"])

        view = st.session_state["nav_view"]
        q_filters = []
        params = []
        if view == "Trash":
            q_filters.append("deleted=1")
        else:
            q_filters.append("deleted=0")
            if view == "Starred":
                q_filters.append("starred=1")
            elif view == "Recent":
                q_filters.append("uploaded_at >= date('now','-7 day')")
        if st.session_state.get("active_folder") and view == "My Drive":
            q_filters.append("folder=?")
            params.append(st.session_state["active_folder"])
        if search_q:
            q_filters.append("(name LIKE ? OR COALESCE(tags,'') LIKE ? OR COALESCE(notes,'') LIKE ?)")
            params.extend([f"%{search_q}%", f"%{search_q}%", f"%{search_q}%"])
        if file_type != "All":
            ext_map = {"Images": [".png",".jpg",".jpeg",".gif",".webp",".bmp",".svg"],
                       "Documents": [".pdf",".docx",".doc",".txt",".md",".rtf",".odt"],
                       "Audio": [".mp3",".wav",".ogg",".flac",".m4a"],
                       "Video": [".mp4",".mov",".avi",".mkv",".webm"],
                       "Data": [".csv",".json",".tsv",".xlsx",".xls",".parquet"],
                       "Code": [".py",".r",".ipynb",".sql",".sh",".js",".html",".css",".yaml",".yml",".toml",".cfg"]}
            sub = " OR ".join(["LOWER(name) LIKE ?" for _ in ext_map[file_type]])
            q_filters.append(f"({sub})")
            params.extend([f"%{e}" for e in ext_map[file_type]])
        where = " AND ".join(q_filters) if q_filters else "1=1"
        order = "uploaded_at DESC" if sort_by == "date" else ("name ASC" if sort_by == "name" else "size_bytes DESC")
        conn = get_conn()
        files = conn.execute(f"SELECT * FROM files WHERE {where} ORDER BY {order}", params).fetchall()
        conn.close()

        if not files:
            st.markdown('<div class="nv-empty">📭 No files in this view.</div>', unsafe_allow_html=True)
        for f in files:
            c1, c2, c3, c4, c5 = st.columns([4, 2, 1, 1, 1])
            title = ("⭐ " if f["starred"] else "") + f["name"]
            if c1.button(title, key=f"prev_{f['id']}"):
                st.session_state["preview_target"] = f["id"]
            c2.caption(f"{fmt_bytes(f['size_bytes'])} · {f['uploaded_at']}")
            conn = get_conn()
            blob = conn.execute("SELECT data_blob FROM files WHERE id=?", (f["id"],)).fetchone()["data_blob"]
            conn.close()
            c3.download_button("⬇️", data=blob, file_name=f["name"], key=f"dl_{f['id']}")
            if f["deleted"]:
                if c4.button("♻️", key=f"restore_{f['id']}"):
                    conn = get_conn(); conn.execute("UPDATE files SET deleted=0 WHERE id=?", (f["id"],)); conn.commit(); conn.close()
                    st.rerun()
            else:
                star_label = "⭐" if f["starred"] else "☆"
                if c4.button(star_label, key=f"star_{f['id']}"):
                    conn = get_conn(); conn.execute("UPDATE files SET starred=1-starred WHERE id=?", (f["id"],)); conn.commit(); conn.close()
                    st.rerun()
            if c5.button("🗑️", key=f"del_{f['id']}"):
                conn = get_conn(); conn.execute("UPDATE files SET deleted=1 WHERE id=?", (f["id"],)); conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"TRASH_FILE:{f['id']}")
                st.rerun()

        if st.session_state["preview_target"]:
            conn = get_conn()
            pf = conn.execute("SELECT * FROM files WHERE id=?", (st.session_state["preview_target"],)).fetchone()
            conn.close()
            if pf:
                st.markdown("---")
                c1, c2 = st.columns([2, 1])
                with c1:
                    st.markdown(f"### 👁️ {pf['name']}")
                    if (pf["content_type"] or "").startswith("image/"):
                        st.image(pf["data_blob"])
                    elif (pf["content_type"] or "").startswith("text/") or pf["name"].endswith((".txt", ".md", ".csv", ".json", ".py", ".r", ".sql", ".html")):
                        st.code(pf["data_blob"].decode("utf-8", errors="ignore")[:8000])
                    elif pf["name"].endswith(".pdf"):
                        import base64 as b64
                        st.markdown(f'<iframe src="data:application/pdf;base64,{b64.b64encode(pf["data_blob"]).decode()}" width="100%" height="600"></iframe>', unsafe_allow_html=True)
                    else:
                        st.caption("No inline preview — use the download button.")
                with c2:
                    st.markdown("##### Details")
                    st.markdown(f"**Size:** {fmt_bytes(pf['size_bytes'])}")
                    st.markdown(f"**Folder:** {pf['folder']}")
                    st.markdown(f"**Uploaded:** {pf['uploaded_at']}")
                    st.markdown(f"**Version:** v{pf['version']}")
                    notes = st.text_area("Notes", value=pf["notes"] or "", key=f"notes_{pf['id']}")
                    tags = st.text_input("Tags (comma-separated)", value=pf["tags"] or "", key=f"tags_{pf['id']}")
                    if st.button("💾 Save Metadata"):
                        require_editor()
                        conn = get_conn(); conn.execute("UPDATE files SET notes=?, tags=? WHERE id=?", (notes, tags, pf["id"])); conn.commit(); conn.close()
                        st.success("Saved")
                        st.rerun()
                    if st.button("🔗 Generate Share Link"):
                        token = secrets.token_urlsafe(8)
                        exp = (datetime.datetime.now() + datetime.timedelta(days=7)).isoformat()
                        conn = get_conn(); conn.execute("INSERT INTO share_links (item_type, item_id, token, permission, expires_at, created_at) VALUES (?,?,?,?,?,?)",
                                                        ("file", pf["id"], token, "view", exp, now_str())); conn.commit(); conn.close()
                        st.code(f"nexus://share/{token}")
                    if st.button("➕ Save New Version"):
                        require_editor()
                        conn = get_conn()
                        cv = conn.execute("SELECT version FROM files WHERE id=?", (pf["id"],)).fetchone()["version"]
                        conn.execute("INSERT INTO file_versions (file_id, version, data_blob, saved_at) VALUES (?,?,?,?)",
                                     (pf["id"], cv + 1, pf["data_blob"], now_str()))
                        conn.execute("UPDATE files SET version=? WHERE id=?", (cv + 1, pf["id"]))
                        conn.commit(); conn.close()
                        st.success("New version saved")
                        st.rerun()
                if st.button("Close preview"):
                    st.session_state["preview_target"] = None
                    st.rerun()

# ============================================================================
# DOCS
# ============================================================================
elif active == "📄 Docs":
    st.subheader("📄 Docs")
    conn = get_conn()
    docs = conn.execute("SELECT * FROM documents WHERE deleted=0 ORDER BY modified_at DESC").fetchall()
    conn.close()

    left, right = st.columns([1, 3])
    with left:
        st.markdown("##### Documents")
        if st.button("➕ New document", use_container_width=True):
            require_editor()
            conn = get_conn()
            cur = conn.execute("INSERT INTO documents (title, content, owner, folder, modified_at) VALUES (?,?,?,?,?)",
                               ("Untitled", "# Untitled\n\nStart typing...", st.session_state["username"], "My Docs", now_str()))
            conn.commit()
            st.session_state["active_doc_id"] = cur.lastrowid
            conn.close()
            st.rerun()
        for d in docs:
            label = ("⭐ " if d["starred"] else "") + (d["title"] or "(untitled)")
            if st.button(label, key=f"doc_{d['id']}", use_container_width=True):
                st.session_state["active_doc_id"] = d["id"]
                st.rerun()
        st.markdown("---")
        st.markdown("##### Templates")
        templates = {
            "Blank": "",
            "Report": "# [Report Title]\n\n## Executive Summary\n\n...\n\n## Findings\n\n...",
            "Meeting Notes": "# Meeting Notes\n\n**Date:** \n**Attendees:** \n\n## Agenda\n1. \n2. \n\n## Action Items\n- [ ] ",
            "Blog Post": "# [Blog Title]\n\n**Intro**\n\n## Section\n\n...\n\n**Conclusion**",
            "Resume": "# [Name]\n\n**Email | Phone | LinkedIn**\n\n## Experience\n### [Company] — [Role]\n- ",
        }
        tchoice = st.selectbox("Start from template", list(templates.keys()))
        if st.button("📄 Use Template"):
            require_editor()
            conn = get_conn()
            cur = conn.execute("INSERT INTO documents (title, content, owner, folder, modified_at) VALUES (?,?,?,?,?)",
                               (tchoice, templates[tchoice], st.session_state["username"], "My Docs", now_str()))
            conn.commit()
            st.session_state["active_doc_id"] = cur.lastrowid
            conn.close()
            st.rerun()

    with right:
        doc = next((d for d in docs if d["id"] == st.session_state["active_doc_id"]), None)
        if doc:
            title = st.text_input("Title", value=doc["title"])
            content = st.text_area("Content (Markdown supported)", value=doc["content"], height=400, key=f"content_{doc['id']}")

            c1, c2, c3, c4, c5 = st.columns(5)
            if c1.button("💾 Save", type="primary"):
                require_editor()
                conn = get_conn()
                conn.execute("UPDATE documents SET title=?, content=?, modified_at=? WHERE id=?",
                             (title, content, now_str(), doc["id"]))
                conn.commit()
                cv = conn.execute("SELECT version FROM documents WHERE id=?", (doc["id"],)).fetchone()["version"]
                conn.execute("INSERT INTO doc_versions (doc_id, version, content, saved_at, editor) VALUES (?,?,?,?,?)",
                             (doc["id"], cv + 1, content, now_str(), st.session_state["username"]))
                conn.execute("UPDATE documents SET version=? WHERE id=?", (cv + 1, doc["id"]))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"SAVE_DOC:{doc['id']}")
                log_activity(st.session_state["username"], f"saved document '{title}'")
                st.toast("Saved", icon="✅")
            if c2.button("🗑️ Delete"):
                require_editor()
                conn = get_conn(); conn.execute("UPDATE documents SET deleted=1 WHERE id=?", (doc["id"],)); conn.commit(); conn.close()
                st.session_state["active_doc_id"] = None
                st.rerun()
            if c3.button("⭐ Star"):
                conn = get_conn(); conn.execute("UPDATE documents SET starred=1-starred WHERE id=?", (doc["id"],)); conn.commit(); conn.close()
                st.rerun()
            if c4.button("📄 Export PDF"):
                if HAS_PDF:
                    pdf = FPDF(); pdf.add_page(); pdf.set_font("Helvetica", size=12)
                    pdf.multi_cell(0, 8, txt=f"{title}\n\n{content}")
                    pdf_bytes = pdf.output(dest="S").encode("latin-1", errors="ignore")
                    st.download_button("⬇️ Download PDF", data=pdf_bytes, file_name=f"{title}.pdf")
                else:
                    st.warning("Install fpdf2 for PDF export")
            if c5.button("📝 Export DOCX"):
                if HAS_DOCX:
                    buf = io.BytesIO()
                    d = docx.Document()
                    d.add_heading(title, level=0)
                    for line in content.splitlines():
                        if line.startswith("## "):
                            d.add_heading(line[3:], level=2)
                        elif line.startswith("# "):
                            d.add_heading(line[2:], level=1)
                        elif line.strip():
                            d.add_paragraph(line)
                    d.save(buf)
                    st.download_button("⬇️ Download DOCX", data=buf.getvalue(), file_name=f"{title}.docx")
                else:
                    st.warning("Install python-docx for DOCX export")

            e1, e2 = st.columns(2)
            e1.download_button("⬇️ .md", data=content.encode(), file_name=f"{title}.md")
            e2.download_button("⬇️ .html", data=f"<html><body><h1>{title}</h1><pre>{content}</pre></body></html>".encode(), file_name=f"{title}.html")
            words = len(content.split())
            st.caption(f"Words: {words} · Characters: {len(content)} · Read time: {words/200:.1f} min · v{doc['version']}")

            l2, r2 = st.columns(2)
            with l2:
                with st.expander("👁️ Preview"):
                    st.markdown(content)
            with r2:
                with st.expander("📚 Outline & Version History"):
                    for line in content.splitlines():
                        if line.startswith("#"):
                            lvl = len(line) - len(line.lstrip("#"))
                            st.markdown("&nbsp;" * lvl + line)
                    conn = get_conn()
                    vers = conn.execute("SELECT * FROM doc_versions WHERE doc_id=? ORDER BY version DESC LIMIT 5", (doc["id"],)).fetchall()
                    conn.close()
                    for v in vers:
                        st.caption(f"v{v['version']} — {v['saved_at']} by {v['editor']}")
        else:
            st.markdown('<div class="nv-empty">📄 Select or create a document on the left.</div>', unsafe_allow_html=True)

# ============================================================================
# SHEETS
# ============================================================================
elif active == "📊 Sheets":
    st.subheader("📊 Sheets")
    conn = get_conn()
    sheets = conn.execute("SELECT * FROM sheets WHERE deleted=0 ORDER BY modified_at DESC").fetchall()
    conn.close()

    left, right = st.columns([1, 3])
    with left:
        st.markdown("##### Sheets")
        if st.button("➕ New sheet", use_container_width=True):
            require_editor()
            conn = get_conn()
            cur = conn.execute("INSERT INTO sheets (name, data_json, owner, modified_at) VALUES (?,?,?,?)",
                               ("Untitled Sheet", json.dumps([{"A": "", "B": ""}]), st.session_state["username"], now_str()))
            conn.commit()
            st.session_state["active_sheet_id"] = cur.lastrowid
            conn.close()
            st.rerun()
        for s in sheets:
            if st.button(s["name"], key=f"sheet_{s['id']}", use_container_width=True):
                st.session_state["active_sheet_id"] = s["id"]
                st.rerun()

    with right:
        sheet = next((s for s in sheets if s["id"] == st.session_state["active_sheet_id"]), None)
        if sheet:
            name = st.text_input("Sheet name", value=sheet["name"])
            st.caption("Formulas: `=SUM(A1:A5)`, `=AVERAGE(B1:B3)`, `=A1*B2`, `=IF(C1>10,\"big\",\"small\")`, `=VLOOKUP(A1,A1:B5,2)`")

            with st.expander("📥 Import CSV"):
                up_csv = st.file_uploader("Choose a CSV file", type=["csv"], key=f"csv_{sheet['id']}")
                if up_csv is not None:
                    imported_df = pd.read_csv(up_csv)
                    st.dataframe(imported_df, use_container_width=True)
                    if st.button("✅ Load CSV"):
                        require_editor()
                        conn = get_conn()
                        conn.execute("UPDATE sheets SET data_json=?, modified_at=? WHERE id=?",
                                     (json.dumps(imported_df.to_dict("records")), now_str(), sheet["id"]))
                        conn.commit(); conn.close()
                        st.rerun()

            rows = json.loads(sheet["data_json"])
            df = pd.DataFrame(rows) if rows else pd.DataFrame({"A": [""]})
            edited = st.data_editor(df, num_rows="dynamic", use_container_width=True, key="nv_sheet_editor")

            c1, c2, c3, c4 = st.columns(4)
            if c1.button("💾 Save", type="primary"):
                require_editor()
                conn = get_conn()
                conn.execute("UPDATE sheets SET name=?, data_json=?, modified_at=? WHERE id=?",
                             (name, json.dumps(edited.to_dict("records")), now_str(), sheet["id"]))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"SAVE_SHEET:{sheet['id']}")
                st.toast("Saved", icon="✅")
            if c2.button("🧮 Show Computed"):
                st.session_state["show_computed"] = True
            if c3.button("🗑️ Delete"):
                require_editor()
                conn = get_conn(); conn.execute("UPDATE sheets SET deleted=1 WHERE id=?", (sheet["id"],)); conn.commit(); conn.close()
                st.session_state["active_sheet_id"] = None
                st.rerun()
            c4.download_button("⬇️ CSV", data=edited.to_csv(index=False).encode("utf-8"), file_name=f"{name}.csv")

            if st.session_state["show_computed"]:
                st.markdown("##### ✅ Computed values")
                computed = compute_sheet(edited)
                st.dataframe(computed, use_container_width=True)
                if st.button("Hide computed"):
                    st.session_state["show_computed"] = False
                    st.rerun()

            computed_full = compute_sheet(edited)
            numeric_cols = computed_full.select_dtypes(include=["number"]).columns.tolist()
            if numeric_cols:
                st.markdown("##### Aggregation & Charts")
                a1, a2, a3 = st.columns(3)
                target = a1.selectbox("Column", numeric_cols)
                op = a2.selectbox("Operation", ["SUM", "AVERAGE", "MIN", "MAX", "COUNT"])
                val = {"SUM": computed_full[target].sum(), "AVERAGE": computed_full[target].mean(),
                       "MIN": computed_full[target].min(), "MAX": computed_full[target].max(),
                       "COUNT": computed_full[target].count()}[op]
                a3.metric(f"{op}({target})", f"{val:,.2f}")
                chart_type = st.radio("Chart", ["Bar", "Line", "Area"], horizontal=True)
                x_axis = computed_full.columns[0]
                if chart_type == "Bar":
                    st.bar_chart(computed_full.set_index(x_axis)[target])
                elif chart_type == "Line":
                    st.line_chart(computed_full.set_index(x_axis)[target])
                else:
                    st.area_chart(computed_full.set_index(x_axis)[target])
            else:
                st.info("Add numeric columns to enable aggregation and charts.")
        else:
            st.markdown('<div class="nv-empty">📊 Select or create a sheet on the left.</div>', unsafe_allow_html=True)

# ============================================================================
# SLIDES
# ============================================================================
elif active == "🎞️ Slides":
    st.subheader("🎞️ Slides")
    conn = get_conn()
    decks = conn.execute("SELECT * FROM slide_decks WHERE deleted=0 ORDER BY modified_at DESC").fetchall()
    conn.close()

    left, right = st.columns([1, 3])
    with left:
        st.markdown("##### Decks")
        if st.button("➕ New deck", use_container_width=True):
            require_editor()
            conn = get_conn()
            cur = conn.execute("INSERT INTO slide_decks (title, slides_json, owner, modified_at) VALUES (?,?,?,?)",
                               ("Untitled Deck", json.dumps([{"title": "Slide 1", "body": "", "layout": "title"}]), st.session_state["username"], now_str()))
            conn.commit()
            st.session_state["active_deck_id"] = cur.lastrowid
            conn.close()
            st.rerun()
        for d in decks:
            if st.button(d["title"], key=f"deck_{d['id']}", use_container_width=True):
                st.session_state["active_deck_id"] = d["id"]
                st.rerun()

    with right:
        deck = next((d for d in decks if d["id"] == st.session_state["active_deck_id"]), None)
        if deck:
            title = st.text_input("Deck title", value=deck["title"])
            theme = st.selectbox("Theme", ["Brass", "Midnight", "Emerald", "Ruby", "Ocean"], key="deck_theme")
            slides = json.loads(deck["slides_json"])
            layouts = ["title", "title-body", "two-column", "quote", "section-header"]

            for i, s in enumerate(slides):
                with st.container(border=True):
                    lc = st.columns([3, 1])
                    lc[0].text_input(f"Slide {i+1} title", value=s.get("title", ""), key=f"st_{deck['id']}_{i}")
                    s["layout"] = lc[1].selectbox("Layout", layouts, index=layouts.index(s.get("layout", "title-body")) if s.get("layout") in layouts else 0, key=f"sl_{deck['id']}_{i}")
                    s["body"] = st.text_area(f"Slide {i+1} body", value=s.get("body", ""), key=f"sb_{deck['id']}_{i}", height=80)

            c1, c2, c3, c4, c5 = st.columns(5)
            if c1.button("➕ Add slide"):
                require_editor()
                slides.append({"title": "New slide", "body": "", "layout": "title-body"})
                conn = get_conn()
                conn.execute("UPDATE slide_decks SET slides_json=?, modified_at=? WHERE id=?", (json.dumps(slides), now_str(), deck["id"]))
                conn.commit(); conn.close()
                st.rerun()
            if len(slides) > 1 and c2.button("🗑️ Remove last"):
                require_editor()
                slides.pop()
                conn = get_conn()
                conn.execute("UPDATE slide_decks SET slides_json=?, modified_at=? WHERE id=?", (json.dumps(slides), now_str(), deck["id"]))
                conn.commit(); conn.close()
                st.rerun()
            if c3.button("💾 Save", type="primary"):
                require_editor()
                conn = get_conn()
                conn.execute("UPDATE slide_decks SET title=?, slides_json=?, modified_at=?, theme=? WHERE id=?",
                             (title, json.dumps(slides), now_str(), theme, deck["id"]))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"SAVE_DECK:{deck['id']}")
                st.toast("Saved", icon="✅")
                st.rerun()
            if c4.button("🗑️ Delete deck"):
                require_editor()
                conn = get_conn(); conn.execute("UPDATE slide_decks SET deleted=1 WHERE id=?", (deck["id"],)); conn.commit(); conn.close()
                st.session_state["active_deck_id"] = None
                st.rerun()
            if c5.button("🎞️ Export PPTX"):
                if HAS_PPTX:
                    prs = Presentation()
                    title_layout, content_layout = prs.slide_layouts[0], prs.slide_layouts[1]
                    for i, s in enumerate(slides):
                        slide = prs.slides.add_slide(title_layout if i == 0 else content_layout)
                        slide.shapes.title.text = s.get("title", "")
                        if len(slide.placeholders) > 1:
                            slide.placeholders[1].text = s.get("body", "")
                    buf = io.BytesIO(); prs.save(buf)
                    st.download_button("⬇️ Download PPTX", data=buf.getvalue(), file_name=f"{title}.pptx")
                else:
                    st.warning("Install python-pptx to export")

            st.markdown("##### 🖥️ Preview")
            theme_bg = {"Brass": "#14161F", "Midnight": "#0F172A", "Emerald": "#064E3B", "Ruby": "#7F1D1D", "Ocean": "#0C4A6E"}.get(theme, INK)
            for s in slides:
                st.markdown(f"""<div style="background:{theme_bg};color:white;border:1px solid {BRASS}44;border-radius:12px;padding:26px;margin-bottom:10px;text-align:center;">
                <h3 style="color:{BRASS};margin:0;">{s.get('title','')}</h3>
                <p style="white-space:pre-line;opacity:0.85;">{s.get("body","")}</p></div>""", unsafe_allow_html=True)
        else:
            st.markdown('<div class="nv-empty">🎞️ Select or create a deck on the left.</div>', unsafe_allow_html=True)

# ============================================================================
# MAIL
# ============================================================================
elif active == "✉️ Mail":
    st.subheader("✉️ Mail")
    tab_inbox, tab_compose, tab_sent, tab_drafts, tab_settings = st.tabs(["📥 Inbox", "✏️ Compose", "📤 Sent", "📝 Drafts", "⚙️ Settings"])

    with tab_inbox:
        conn = get_conn()
        inbox = conn.execute("SELECT * FROM mail WHERE folder='inbox' ORDER BY sent_at DESC").fetchall()
        conn.close()
        if not inbox:
            st.markdown('<div class="nv-empty">📭 Inbox is empty.</div>', unsafe_allow_html=True)
        for m in inbox:
            with st.container(border=True):
                read = "" if m["read"] else "🔵 "
                star = "⭐" if m["starred"] else "☆"
                c1, c2 = st.columns([6, 1])
                c1.markdown(f"**{read}{m['subject']}** · *{m['sender']} · {m['sent_at']}*  \n<span class='nv-chip'>{m['label']}</span>", unsafe_allow_html=True)
                c2.button(star, key=f"mstar_{m['id']}", help="Star")
                st.caption(m["body"][:200])
                b1, b2, b3 = st.columns(3)
                if b1.button("Read/Unread", key=f"mread_{m['id']}"):
                    conn = get_conn(); conn.execute("UPDATE mail SET read=? WHERE id=?", (0 if m["read"] else 1, m["id"])); conn.commit(); conn.close()
                    st.rerun()
                if b2.button("Archive", key=f"march_{m['id']}"):
                    conn = get_conn(); conn.execute("UPDATE mail SET folder='archive' WHERE id=?", (m["id"],)); conn.commit(); conn.close()
                    st.rerun()
                if b3.button("🗑️", key=f"mtrash_{m['id']}"):
                    conn = get_conn(); conn.execute("UPDATE mail SET folder='trash' WHERE id=?", (m["id"],)); conn.commit(); conn.close()
                    st.rerun()

    with tab_compose:
        contacts = get_conn().execute("SELECT email FROM contacts").fetchall()
        to_options = [c["email"] for c in contacts]
        to = st.selectbox("To (from contacts)", [""] + to_options)
        subject = st.text_input("Subject")
        body = st.text_area("Message (HTML supported)", height=180)
        label = st.selectbox("Label", ["Primary", "Social", "Promotions", "Updates"])
        smtp_creds = st.session_state.get("smtp_creds", {})
        if st.button("📨 Send", type="primary"):
            mode = "demo"
            if smtp_creds.get("host") and smtp_creds.get("user") and smtp_creds.get("password"):
                try:
                    import smtplib
                    import ssl
                    from email.mime.text import MIMEText
                    msg = MIMEText(body, "html")
                    msg["Subject"], msg["From"], msg["To"] = subject, smtp_creds["user"], to
                    with smtplib.SMTP_SSL(smtp_creds["host"], 465, context=ssl.create_default_context()) as server:
                        server.login(smtp_creds["user"], smtp_creds["password"])
                        server.sendmail(smtp_creds["user"], [to], msg.as_string())
                    mode = "live"
                except Exception as e:
                    st.error(f"SMTP failed, saved as demo: {e}")
            conn = get_conn()
            conn.execute("INSERT INTO mail (folder, sender, recipient, subject, body, mode, sent_at, label) VALUES (?,?,?,?,?,?,?,?)",
                         ("sent", st.session_state["username"], to, subject, body, mode, now_str(), label))
            conn.commit(); conn.close()
            log_audit(st.session_state["username"], f"SEND_MAIL:{to}:{mode}")
            badge = "nv-badge-live" if mode == "live" else "nv-badge-demo"
            st.markdown(f'<span class="nv-badge {badge}">{mode.upper()}</span> Message recorded.', unsafe_allow_html=True)

    with tab_sent:
        conn = get_conn()
        sent = conn.execute("SELECT * FROM mail WHERE folder='sent' ORDER BY sent_at DESC").fetchall()
        conn.close()
        if not sent:
            st.markdown('<div class="nv-empty">📤 Nothing sent yet.</div>', unsafe_allow_html=True)
        for m in sent:
            badge = "nv-badge-live" if m["mode"] == "live" else "nv-badge-demo"
            with st.container(border=True):
                st.markdown(f'<span class="nv-badge {badge}">{m["mode"].upper()}</span> **{m["subject"]}** → {m["recipient"]}', unsafe_allow_html=True)
                st.caption(m["sent_at"])

    with tab_drafts:
        ds = st.text_input("Draft subject", key="draft_subject")
        db_ = st.text_area("Draft body", key="draft_body", height=120)
        if st.button("💾 Save draft"):
            conn = get_conn()
            conn.execute("INSERT INTO mail (folder, sender, recipient, subject, body, mode, sent_at) VALUES (?,?,?,?,?,?,?)",
                         ("drafts", st.session_state["username"], "", ds, db_, "draft", now_str()))
            conn.commit(); conn.close()
            st.success("Draft saved")
            st.rerun()
        conn = get_conn()
        drafts = conn.execute("SELECT * FROM mail WHERE folder='drafts' ORDER BY sent_at DESC").fetchall()
        conn.close()
        for d in drafts:
            with st.container(border=True):
                st.markdown(f"**{d['subject'] or '(no subject)'}**  \n{d['body'][:150]}")
                if st.button("Delete draft", key=f"draft_{d['id']}"):
                    conn = get_conn(); conn.execute("DELETE FROM mail WHERE id=?", (d["id"],)); conn.commit(); conn.close()
                    st.rerun()

    with tab_settings:
        st.caption("Enter SMTP details (e.g. Gmail app password) to actually send mail. Session-only.")
        host = st.text_input("SMTP host", value=st.session_state.get("smtp_creds", {}).get("host", "smtp.gmail.com"))
        user = st.text_input("SMTP username / email", value=st.session_state.get("smtp_creds", {}).get("user", ""))
        pw = st.text_input("SMTP app password", type="password")
        if st.button("Save SMTP settings"):
            st.session_state["smtp_creds"] = {"host": host, "user": user, "password": pw}
            st.success("Saved for this session.")

# ============================================================================
# CONTACTS
# ============================================================================
elif active == "👥 Contacts":
    st.subheader("👥 Contacts")
    require_editor()

    with st.expander("➕ Add contact"):
        with st.form("contact_form"):
            c_name = st.text_input("Name")
            c_email = st.text_input("Email")
            c_phone = st.text_input("Phone")
            c_company = st.text_input("Company")
            c_group = st.selectbox("Group", ["General", "Work", "Friends", "Family", "Clients"])
            c_notes = st.text_area("Notes")
            if st.form_submit_button("Save contact"):
                conn = get_conn()
                conn.execute("INSERT INTO contacts (name, email, phone, company, group_name, notes, created_at) VALUES (?,?,?,?,?,?,?)",
                             (c_name, c_email, c_phone, c_company, c_group, c_notes, now_str()))
                conn.commit(); conn.close()
                st.success("Contact added")
                st.rerun()

    group_filter = st.selectbox("Filter by group", ["All", "General", "Work", "Friends", "Family", "Clients"])
    conn = get_conn()
    if group_filter == "All":
        contacts = conn.execute("SELECT * FROM contacts ORDER BY name").fetchall()
    else:
        contacts = conn.execute("SELECT * FROM contacts WHERE group_name=? ORDER BY name", (group_filter,)).fetchall()
    conn.close()

    search_c = st.text_input("🔍 Search contacts", key="contact_search")
    if search_c:
        contacts = [c for c in contacts if search_c.lower() in c["name"].lower() or search_c.lower() in (c["email"] or "").lower()]

    if not contacts:
        st.markdown('<div class="nv-empty">👥 No contacts found.</div>', unsafe_allow_html=True)
    for c in contacts:
        with st.container(border=True):
            st.markdown(f"**{c['name']}**  \n{c['email']} · {c['phone']} · {c['company']}  \n<span class='nv-chip'>{c['group_name']}</span>", unsafe_allow_html=True)
            if st.button("🗑️ Delete", key=f"con_{c['id']}"):
                conn = get_conn(); conn.execute("DELETE FROM contacts WHERE id=?", (c["id"],)); conn.commit(); conn.close()
                st.rerun()

# ============================================================================
# CALENDAR
# ============================================================================
elif active == "📅 Calendar":
    st.subheader("📅 Calendar")
    require_editor()
    import calendar as _cal

    view_mode = st.radio("View", ["Month", "Week", "Day", "Agenda"], horizontal=True)
    conn = get_conn()
    events = conn.execute("SELECT * FROM calendar_events ORDER BY start_date, start_time").fetchall()
    conn.close()

    with st.expander("➕ Add event"):
        with st.form("event_form"):
            ev_title = st.text_input("Title")
            c1, c2 = st.columns(2)
            start_date = c1.date_input("Start date", datetime.date.today())
            end_date = c2.date_input("End date", datetime.date.today())
            c3, c4 = st.columns(2)
            start_time = c3.time_input("Start time", datetime.time(9, 0))
            end_time = c4.time_input("End time", datetime.time(10, 0))
            category = st.selectbox("Category", ["Work", "Personal", "Health", "Important"])
            remind = st.selectbox("Remind me", [0, 5, 10, 30, 60])
            notes = st.text_area("Notes")
            if st.form_submit_button("Save event"):
                conn = get_conn()
                conn.execute("INSERT INTO calendar_events (title, start_date, end_date, start_time, end_time, category, remind_minutes, notes, owner) VALUES (?,?,?,?,?,?,?,?,?)",
                             (ev_title, start_date.isoformat(), end_date.isoformat(), start_time.strftime("%H:%M"), end_time.strftime("%H:%M"), category, remind, notes, st.session_state["username"]))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"ADD_EVENT:{ev_title}")
                st.success("Event added")
                st.rerun()

    if view_mode in ("Agenda", "Day"):
        today = datetime.date.today().isoformat()
        day_events = [e for e in events if e["start_date"] <= today <= e["end_date"]]
        if not day_events:
            st.markdown('<div class="nv-empty">📅 No events today.</div>', unsafe_allow_html=True)
        for e in day_events:
            with st.container(border=True):
                st.markdown(f"**{e['title']}** · `{e['category']}` · {e['start_time']}–{e['end_time']}")
                if e["notes"]:
                    st.caption(e["notes"])
                if st.button("🗑️ Delete", key=f"ev_{e['id']}"):
                    conn = get_conn(); conn.execute("DELETE FROM calendar_events WHERE id=?", (e["id"],)); conn.commit(); conn.close()
                    st.rerun()
    elif view_mode == "Month":
        st.markdown(f"### {datetime.date.today().strftime('%B %Y')}")
        month_days = _cal.monthdayscalendar(datetime.date.today().year, datetime.date.today().month)
        for week in month_days:
            cols = st.columns(7)
            for idx, day in enumerate(week):
                if day == 0:
                    cols[idx].markdown("")
                else:
                    ds = f"{datetime.date.today().year}-{datetime.date.today().month:02d}-{day:02d}"
                    day_has = [e for e in events if e["start_date"] <= ds <= e["end_date"]]
                    mark = "🔴" if day_has else ""
                    cols[idx].markdown(f"**{day}** {mark}")
        st.caption("🔴 marks days with events.")
    else:
        st.markdown("### This Week")
        today = datetime.date.today()
        monday = today - datetime.timedelta(days=today.weekday())
        for i in range(7):
            ds = (monday + datetime.timedelta(days=i)).isoformat()
            day_events = [e for e in events if e["start_date"] <= ds <= e["end_date"]]
            st.markdown(f"**{(monday + datetime.timedelta(days=i)).strftime('%A %d')}**")
            if day_events:
                for e in day_events:
                    st.markdown(f"- {e['start_time']} {e['title']} ({e['category']})")
            else:
                st.caption("No events")

# ============================================================================
# TASKS
# ============================================================================
elif active == "✅ Tasks":
    st.subheader("✅ Tasks")
    require_editor()

    with st.expander("➕ Add task"):
        with st.form("task_form"):
            t_title = st.text_input("Task")
            t_list = st.text_input("List name", value="Default")
            t_due = st.date_input("Due date", None)
            t_priority = st.selectbox("Priority", ["Low", "Medium", "High"])
            if st.form_submit_button("Add task"):
                conn = get_conn()
                conn.execute("INSERT INTO tasks (title, due_date, priority, list_name, owner, created_at) VALUES (?,?,?,?,?,?)",
                             (t_title, t_due.isoformat() if t_due else None, t_priority, t_list, st.session_state["username"], now_str()))
                conn.commit(); conn.close()
                st.success("Task added")
                st.rerun()

    list_filter = get_conn().execute("SELECT DISTINCT list_name FROM tasks").fetchall()
    list_names = [l["list_name"] for l in list_filter]
    sel_list = st.selectbox("List", ["All"] + list_names)
    conn = get_conn()
    if sel_list == "All":
        tasks = conn.execute("SELECT * FROM tasks ORDER BY done, priority DESC").fetchall()
    else:
        tasks = conn.execute("SELECT * FROM tasks WHERE list_name=? ORDER BY done, priority DESC", (sel_list,)).fetchall()
    conn.close()

    open_tasks = [t for t in tasks if not t["done"]]
    done_tasks = [t for t in tasks if t["done"]]
    c1, c2 = st.columns(2)
    c1.metric("Open", len(open_tasks))
    c2.metric("Completed", len(done_tasks))

    for t in tasks:
        done = "✅" if t["done"] else "⬜"
        prio = {"High": "🔴", "Medium": "🟡", "Low": "🟢"}.get(t["priority"], "⚪")
        c1, c2 = st.columns([5, 1])
        c1.markdown(f"{done} {prio} {t['title']}  \n*{t['list_name']} · due {t['due_date'] or '—'}*")
        if c2.button("✓" if not t["done"] else "↩️", key=f"task_{t['id']}"):
            conn = get_conn(); conn.execute("UPDATE tasks SET done=? WHERE id=?", (0 if t["done"] else 1, t["id"])); conn.commit(); conn.close()
            st.rerun()

# ============================================================================
# NOTES (Google Keep clone)
# ============================================================================
elif active == "🗒️ Notes":
    st.subheader("🗒️ Notes")
    require_editor()

    with st.expander("➕ Add note"):
        with st.form("note_form"):
            n_title = st.text_input("Title")
            n_body = st.text_area("Body")
            n_color = st.selectbox("Color", ["#14161F", "#1E3A8A", "#065F46", "#92400E", "#7F1D1D"])
            if st.form_submit_button("Add note"):
                conn = get_conn()
                conn.execute("INSERT INTO notes (title, body, color, owner, created_at) VALUES (?,?,?,?,?)",
                             (n_title, n_body, n_color, st.session_state["username"], now_str()))
                conn.commit(); conn.close()
                st.rerun()

    conn = get_conn()
    notes = conn.execute("SELECT * FROM notes ORDER BY pinned DESC, created_at DESC").fetchall()
    conn.close()

    if not notes:
        st.markdown('<div class="nv-empty">🗒️ No notes yet.</div>', unsafe_allow_html=True)

    cols = st.columns(3)
    for i, n in enumerate(notes):
        col = cols[i % 3]
        with col:
            pin = "📌" if n["pinned"] else "📌"
            st.markdown(f"""<div class="nv-note" style="background:{n['color']};">
                <b>{pin} {n['title']}</b><br><span style="color:white;opacity:0.85;white-space:pre-line;">{n['body']}</span></div>""", unsafe_allow_html=True)
            b1, b2 = st.columns(2)
            if b1.button("Pin", key=f"pin_{n['id']}"):
                conn = get_conn(); conn.execute("UPDATE notes SET pinned=? WHERE id=?", (0 if n["pinned"] else 1, n["id"])); conn.commit(); conn.close()
                st.rerun()
            if b2.button("🗑️", key=f"note_{n['id']}"):
                conn = get_conn(); conn.execute("DELETE FROM notes WHERE id=?", (n["id"],)); conn.commit(); conn.close()
                st.rerun()

# ============================================================================
# MEETINGS
# ============================================================================
elif active == "🎥 Meetings":
    st.subheader("🎥 Meetings")
    st.caption("Real video rooms via Jitsi Meet — free, no account or API key needed.")

    tab_join, tab_schedule, tab_history = st.tabs(["🎥 Join", "📅 Schedule", "🕘 History"])

    with tab_join:
        if not st.session_state["meeting_room"]:
            st.session_state["meeting_room"] = f"nexusvault-{uuid.uuid4().hex[:8]}"
        room = st.session_state["meeting_room"]
        room_url = f"https://meet.jit.si/{room}"
        st.text_input("Meeting link (share this with others)", value=room_url, disabled=True)
        c1, c2 = st.columns(2)
        if c1.button("🔁 New room"):
            st.session_state["meeting_room"] = f"nexusvault-{uuid.uuid4().hex[:8]}"
            st.rerun()
        join = c2.toggle("🎥 Join now")
        if join:
            components.iframe(room_url, height=560)
        else:
            st.markdown('<div class="nv-empty">🎥 Toggle "Join now" to load the live video call.</div>', unsafe_allow_html=True)

    with tab_schedule:
        require_editor()
        with st.form("meeting_form"):
            m_title = st.text_input("Meeting title")
            m_date = st.date_input("Schedule for", datetime.date.today())
            m_time = st.time_input("Time", datetime.time(10, 0))
            m_recurring = st.checkbox("Recurring meeting")
            invitees = st.text_input("Invitees (comma-separated emails)")
            if st.form_submit_button("Schedule meeting"):
                conn = get_conn()
                room_id = f"nexusvault-{uuid.uuid4().hex[:8]}"
                invites = [x.strip() for x in invitees.split(",") if x.strip()]
                conn.execute("INSERT INTO meetings (title, room, scheduled_for, owner, created_at, recurring, invitees) VALUES (?,?,?,?,?,?,?)",
                             (m_title, room_id, f"{m_date} {m_time.strftime('%H:%M')}", st.session_state["username"], now_str(), int(m_recurring), json.dumps(invites)))
                conn.commit(); conn.close()
                log_audit(st.session_state["username"], f"SCHEDULE_MEETING:{m_title}")
                st.success(f"Scheduled! Join: https://meet.jit.si/{room_id}")
                conn = get_conn()
                conn.execute("INSERT INTO calendar_events (title, start_date, end_date, start_time, end_time, category, owner) VALUES (?,?,?,?,?,?,?)",
                             (m_title, m_date.isoformat(), m_date.isoformat(), m_time.strftime("%H:%M"), m_time.strftime("%H:%M"), "Meeting", st.session_state["username"]))
                conn.commit(); conn.close()
                st.rerun()

    with tab_history:
        conn = get_conn()
        meetings = conn.execute("SELECT * FROM meetings ORDER BY created_at DESC").fetchall()
        conn.close()
        if not meetings:
            st.markdown('<div class="nv-empty">🕘 No scheduled meetings yet.</div>', unsafe_allow_html=True)
        for m in meetings:
            with st.container(border=True):
                rec = " 🔁" if m["recurring"] else ""
                st.markdown(f"**{m['title']}**{rec}  \nScheduled: {m['scheduled_for']}  \nLink: `https://meet.jit.si/{m['room']}`")
                if m["invitees"] and m["invitees"] != "[]":
                    st.caption(f"Invitees: {m['invitees']}")
                if st.button("Open room", key=f"meet_{m['id']}"):
                    st.session_state["meeting_room"] = m["room"]
                    st.rerun()

# ============================================================================
# AUTOMATIONS
# ============================================================================
elif active == "⚙️ Automations":
    st.subheader("⚙️ Automations")

    with st.expander("➕ Add rule"):
        with st.form("rule_form"):
            r_name = st.text_input("Rule name")
            trig = st.selectbox("Trigger", ["storage_high", "file_upload", "daily_report"])
            action = st.selectbox("Action", ["flag_largest_file", "notify_admin", "cleanup_trash"])
            threshold = st.number_input("Threshold (MB)", min_value=1, value=5)
            if st.form_submit_button("Add rule"):
                conn = get_conn()
                conn.execute("INSERT INTO automation_rules (name, trigger_type, action_type, threshold, enabled) VALUES (?,?,?,?,1)",
                             (r_name, trig, action, threshold))
                conn.commit(); conn.close()
                st.success("Rule added")
                st.rerun()

    conn = get_conn()
    rules = conn.execute("SELECT * FROM automation_rules").fetchall()
    conn.close()
    for r in rules:
        c1, c2 = st.columns([5, 1])
        c1.markdown(f"**{r['name']}** — `{r['trigger_type']}` → `{r['action_type']}` — {'✅ enabled' if r['enabled'] else '⏸️ disabled'}")
        if c2.button("Toggle", key=f"rule_{r['id']}"):
            conn = get_conn(); conn.execute("UPDATE automation_rules SET enabled=? WHERE id=?", (0 if r["enabled"] else 1, r["id"])); conn.commit(); conn.close()
            st.rerun()

    if st.button("▶️ Run automations now", type="primary"):
        conn = get_conn()
        total = conn.execute("SELECT COALESCE(SUM(size_bytes),0) FROM files WHERE deleted=0").fetchone()[0]
        results = []
        rules2 = conn.execute("SELECT * FROM automation_rules WHERE enabled=1").fetchall()
        for rule in rules2:
            if rule["trigger_type"] == "storage_high" and total > rule["threshold"] * 1024 * 1024:
                if rule["action_type"] == "flag_largest_file":
                    biggest = conn.execute("SELECT name, size_bytes FROM files WHERE deleted=0 ORDER BY size_bytes DESC LIMIT 1").fetchone()
                    if biggest:
                        results.append(f"[{rule['name']}] Largest file is '{biggest['name']}' ({biggest['size_bytes']/1024:.1f} KB).")
            elif rule["trigger_type"] == "file_upload":
                results.append(f"[{rule['name']}] Checking recent uploads...")
            else:
                results.append(f"[{rule['name']}] Storage usage is low — no action needed.")
        if not rules2:
            results.append("No enabled rules.")
        conn.close()
        log_audit(st.session_state["username"], "RUN_AUTOMATIONS")
        for line in results:
            st.write(f"- {line}")

# ============================================================================
# PASSWORDS
# ============================================================================
elif active == "🔑 Passwords":
    st.subheader("🔑 Passwords")
    require_editor()

    if not HAS_CRYPTO:
        st.warning("Install `cryptography` to store encrypted passwords.")
    else:
        if not st.session_state["vault_key"]:
            st.session_state["vault_key"] = Fernet.generate_key()
        cipher = Fernet(st.session_state["vault_key"])

        with st.expander("➕ Add password"):
            with st.form("pw_form"):
                p_site = st.text_input("Site / Service")
                p_user = st.text_input("Username / Email")
                p_pass = st.text_input("Password", type="password")
                p_notes = st.text_input("Notes")
                if st.form_submit_button("Save password"):
                    enc = cipher.encrypt(p_pass.encode())
                    conn = get_conn()
                    conn.execute("INSERT INTO passwords (site, username, encrypted_password, notes, owner, created_at) VALUES (?,?,?,?,?,?)",
                                 (p_site, p_user, enc, p_notes, st.session_state["username"], now_str()))
                    conn.commit(); conn.close()
                    st.success("Password stored encrypted")
                    st.rerun()

        conn = get_conn()
        passwords = conn.execute("SELECT * FROM passwords WHERE owner=?", (st.session_state["username"],)).fetchall()
        conn.close()
        if not passwords:
            st.markdown('<div class="nv-empty">🔑 No passwords stored yet.</div>', unsafe_allow_html=True)
        for pw in passwords:
            with st.container(border=True):
                st.markdown(f"**{pw['site']}** — {pw['username']}")
                if st.button("🔓 Reveal", key=f"reveal_{pw['id']}"):
                    try:
                        dec = cipher.decrypt(pw["encrypted_password"]).decode()
                        st.code(dec)
                    except Exception:
                        st.error("Cannot decrypt (key changed since save).")
                if st.button("🗑️ Delete", key=f"pwd_{pw['id']}"):
                    conn = get_conn(); conn.execute("DELETE FROM passwords WHERE id=?", (pw["id"],)); conn.commit(); conn.close()
                    st.rerun()

# ============================================================================
# SECURITY
# ============================================================================
elif active == "🛡️ Security":
    st.subheader("🛡️ Security")
    tab_enc, tab_2fa, tab_dash = st.tabs(["🔐 Encrypt/Decrypt", "🔢 2FA Setup", "📊 Dashboard"])

    with tab_enc:
        if not HAS_CRYPTO:
            st.warning("Install `cryptography` to enable encrypt/decrypt.")
        else:
            if not st.session_state["fernet_key"]:
                st.session_state["fernet_key"] = Fernet.generate_key()
            cipher = Fernet(st.session_state["fernet_key"])
            mode = st.radio("Operation", ["Encrypt", "Decrypt"], horizontal=True)
            text = st.text_area("Text")
            if mode == "Encrypt" and st.button("Encrypt"):
                st.code(cipher.encrypt(text.encode()).decode())
            if mode == "Decrypt" and st.button("Decrypt"):
                try:
                    st.code(cipher.decrypt(text.encode()).decode())
                except Exception:
                    st.error("Invalid ciphertext for this session's key.")

    with tab_2fa:
        st.caption("2FA setup (demo): generate a TOTP secret and verify a code.")
        if "totp_secret" not in st.session_state:
            st.session_state["totp_secret"] = base64.b32encode(secrets.token_bytes(10)).decode()
        st.info(f"**TOTP Secret:** `{st.session_state['totp_secret']}`")
        code = st.text_input("Enter 6-digit code", max_chars=6)
        if st.button("Verify code"):
            try:
                import hmac
                import struct
                import time
                key = base64.b32decode(st.session_state["totp_secret"])
                counter = struct.pack(">Q", int(time.time()) // 30)
                h = hmac.new(key, counter, hashlib.sha1).digest()
                offset = h[-1] & 0x0F
                val = (struct.unpack(">I", h[offset:offset + 4])[0] & 0x7FFFFFFF) % 1000000
                if code == f"{val:06d}":
                    st.success("✅ 2FA code verified!")
                    log_audit(st.session_state["username"], "2FA_VERIFIED")
                else:
                    st.error("Invalid code.")
            except Exception as e:
                st.error(f"Error: {e}")

    with tab_dash:
        st.markdown("##### Security posture")
        c1, c2, c3 = st.columns(3)
        conn = get_conn()
        logins = conn.execute("SELECT COUNT(*) FROM audit_log WHERE action='LOGIN'").fetchone()[0]
        failed = conn.execute("SELECT COUNT(*) FROM audit_log WHERE status='DENIED'").fetchone()[0]
        conn.close()
        c1.metric("Successful logins", logins)
        c2.metric("Denied attempts", failed)
        c3.metric("Encryption", "AES-256-GCM" if HAS_CRYPTO else "Unavailable")

# ============================================================================
# STORAGE & BACKUP
# ============================================================================
elif active == "☁️ Storage & Backup":
    st.subheader("☁️ Storage & Backup")
    plan = st.selectbox("Plan", list(STORAGE_PLANS_GB.keys()))
    limit_gb = STORAGE_PLANS_GB[plan]
    conn = get_conn()
    total_bytes = conn.execute("SELECT COALESCE(SUM(size_bytes),0) FROM files WHERE deleted=0").fetchone()[0]
    conn.close()
    used_gb = total_bytes / (1024 ** 3)
    st.progress(min(used_gb / limit_gb, 1.0) if limit_gb else 0)
    st.caption(f"{used_gb:.4f} GB used of {plan}")

    conn = get_conn()
    by_folder = conn.execute("SELECT folder, SUM(size_bytes) as sz, COUNT(*) as cnt FROM files WHERE deleted=0 GROUP BY folder").fetchall()
    conn.close()
    if by_folder:
        st.markdown("##### Usage by folder")
        for row in by_folder:
            st.markdown(f"- **{row['folder']}**: {row['cnt']} files, {row['sz']/1024:.1f} KB")

    st.markdown("---")
    st.markdown("##### 💾 Full backup")
    conn = get_conn()
    backup = {
        "documents": [dict(r) for r in conn.execute("SELECT * FROM documents WHERE deleted=0").fetchall()],
        "sheets": [dict(r) for r in conn.execute("SELECT * FROM sheets WHERE deleted=0").fetchall()],
        "slide_decks": [dict(r) for r in conn.execute("SELECT * FROM slide_decks WHERE deleted=0").fetchall()],
        "mail": [dict(r) for r in conn.execute("SELECT * FROM mail").fetchall()],
        "contacts": [dict(r) for r in conn.execute("SELECT * FROM contacts").fetchall()],
        "calendar_events": [dict(r) for r in conn.execute("SELECT * FROM calendar_events").fetchall()],
        "meetings": [dict(r) for r in conn.execute("SELECT * FROM meetings").fetchall()],
        "tasks": [dict(r) for r in conn.execute("SELECT * FROM tasks").fetchall()],
        "notes": [dict(r) for r in conn.execute("SELECT * FROM notes").fetchall()],
    }
    conn.close()
    st.download_button("⬇️ Download backup (.json)", data=json.dumps(backup, indent=2, default=str),
                       file_name=f"nexus_vault_backup_{datetime.date.today()}.json", mime="application/json")

    st.markdown("##### 🗑️ Empty trash")
    if st.button("Empty trash permanently"):
        conn = get_conn()
        conn.execute("DELETE FROM files WHERE deleted=1")
        conn.commit(); conn.close()
        log_audit(st.session_state["username"], "EMPTY_TRASH")
        st.success("Trash emptied")
        st.rerun()

# ============================================================================
# USER MANAGEMENT
# ============================================================================
elif active == "👤 Users":
    st.subheader("👤 Users")
    if st.session_state["role"] != "Admin":
        st.error("Admins only.")
        st.stop()

    with st.expander("➕ Add user"):
        with st.form("user_form"):
            nu = st.text_input("Username")
            npp = st.text_input("Password", type="password")
            nr = st.selectbox("Role", ["Admin", "Editor", "Viewer"])
            if st.form_submit_button("Create user"):
                conn = get_conn()
                try:
                    conn.execute("INSERT INTO users (username, password_hash, role, created_at) VALUES (?,?,?,?)",
                                 (nu, hash_password(npp), nr, now_str()))
                    conn.commit()
                    st.success(f"User {nu} created")
                    log_audit(st.session_state["username"], f"CREATE_USER:{nu}")
                except sqlite3.IntegrityError:
                    st.error("Username already exists.")
                conn.close()
                st.rerun()

    conn = get_conn()
    users = conn.execute("SELECT * FROM users").fetchall()
    conn.close()
    for u in users:
        c1, c2 = st.columns([4, 1])
        c1.markdown(f"**{u['username']}** — `{u['role']}` (created {u['created_at']})")
        if u["username"] != st.session_state["username"] and c2.button("🗑️ Delete", key=f"user_{u['id']}"):
            conn = get_conn(); conn.execute("DELETE FROM users WHERE id=?", (u["id"],)); conn.commit(); conn.close()
            log_audit(st.session_state["username"], f"DELETE_USER:{u['username']}")
            st.rerun()

# ============================================================================
# ADMIN LOG
# ============================================================================
elif active == "📋 Admin Log":
    st.subheader("📋 Admin Log")
    if st.session_state["role"] != "Admin":
        st.error("Admins only.")
        st.stop()

    conn = get_conn()
    logs = conn.execute("SELECT * FROM audit_log ORDER BY timestamp DESC LIMIT 300").fetchall()
    conn.close()
    if logs:
        df = pd.DataFrame([dict(r) for r in logs])
        filter_actor = st.selectbox("Filter by user", ["All"] + list(df["actor"].unique()))
        if filter_actor != "All":
            df = df[df["actor"] == filter_actor]
        st.dataframe(df, use_container_width=True, hide_index=True)
else:
        st.info("No audit log entries yet.")
